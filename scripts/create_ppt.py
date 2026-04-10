#!/usr/bin/env python3
"""
Generate CoastalCred 10-slide presentation PowerPoint file.
Reads results JSONs and embeds images from results/ directory.
"""

import json
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
ROOT = Path(__file__).resolve().parent.parent
RESULTS = ROOT / "results"
OUTPUT = ROOT / "reports" / "CoastalCred_Presentation.pptx"
OUTPUT.parent.mkdir(parents=True, exist_ok=True)

# ---------------------------------------------------------------------------
# Load data
# ---------------------------------------------------------------------------
with open(RESULTS / "ndvi.json") as f:
    ndvi_data = json.load(f)
with open(RESULTS / "xgboost.json") as f:
    xgb_data = json.load(f)
with open(RESULTS / "unet.json") as f:
    unet_data = json.load(f)
with open(RESULTS / "carbon_predictions.json") as f:
    carbon_data = json.load(f)

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
DARK_BLUE = RGBColor(0x00, 0x77, 0xB6)
GREEN = RGBColor(0x2D, 0x6A, 0x4F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_BLUE_ROW = RGBColor(0xE8, 0xF4, 0xFD)
LIGHT_GREEN = RGBColor(0xD8, 0xF3, 0xDC)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

slide_number_counter = [0]


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------
def add_slide():
    """Add a blank slide and return it."""
    slide_number_counter[0] += 1
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # slide number
    txBox = slide.shapes.add_textbox(
        SLIDE_W - Inches(0.8), SLIDE_H - Inches(0.45), Inches(0.6), Inches(0.35)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(slide_number_counter[0])
    p.font.size = Pt(10)
    p.font.color.rgb = DARK_GRAY
    p.alignment = PP_ALIGN.RIGHT
    return slide


def add_header_bar(slide, title_text):
    """Draw a dark-blue header bar at the top with white title text."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(1.0)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()

    txBox = slide.shapes.add_textbox(
        Inches(0.6), Inches(0.15), SLIDE_W - Inches(1.2), Inches(0.7)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT


def add_green_accent_line(slide, top=Inches(1.05)):
    """Draw a thin green accent line below the header."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), top, SLIDE_W, Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN
    line.line.fill.background()


def add_textbox(slide, text, left, top, width, height,
                font_size=Pt(18), bold=False, color=DARK_GRAY,
                alignment=PP_ALIGN.LEFT, line_spacing=1.3):
    """Add a textbox with given text. Supports newline-separated lines."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = alignment
        p.space_after = Pt(4)
        p.line_spacing = line_spacing
    return txBox


def add_bullet_box(slide, bullets, left, top, width, height,
                   font_size=Pt(16), color=DARK_GRAY, bold_prefix=True):
    """Add bullet points. Each bullet is a string; supports 'Bold: rest' format."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        p.line_spacing = 1.4

        if bold_prefix and ": " in bullet:
            prefix, rest = bullet.split(": ", 1)
            run1 = p.add_run()
            run1.text = "\u2022 " + prefix + ": "
            run1.font.size = font_size
            run1.font.bold = True
            run1.font.color.rgb = color
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = font_size
            run2.font.bold = False
            run2.font.color.rgb = color
        else:
            run = p.add_run()
            run.text = "\u2022 " + bullet
            run.font.size = font_size
            run.font.bold = False
            run.font.color.rgb = color
    return txBox


def add_image(slide, img_path, left, top, width=None, height=None):
    """Embed an image if it exists. Returns True if added, False if not found."""
    img = RESULTS / img_path if not os.path.isabs(str(img_path)) else Path(img_path)
    if not img.exists():
        print(f"  [WARN] Image not found, skipping: {img}")
        return False
    kwargs = {}
    if width:
        kwargs["width"] = width
    if height:
        kwargs["height"] = height
    slide.shapes.add_picture(str(img), left, top, **kwargs)
    return True


def add_table(slide, rows, col_widths, left, top, height_per_row=Inches(0.45),
              header_color=DARK_BLUE, alt_color=LIGHT_BLUE_ROW,
              font_size=Pt(14), header_font_size=Pt(15)):
    """Add a formatted table with alternating row colors."""
    n_rows = len(rows)
    n_cols = len(rows[0]) if rows else 0
    total_width = sum(col_widths)
    total_height = height_per_row * n_rows

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top,
                                          Inches(total_width / 914400),
                                          total_height)
    table = table_shape.table

    # Set column widths
    for j, w in enumerate(col_widths):
        table.columns[j].width = w

    for i, row in enumerate(rows):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(cell_text)

            # Clear default paragraph formatting
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = header_font_size if i == 0 else font_size
                paragraph.font.bold = (i == 0)
                paragraph.font.color.rgb = WHITE if i == 0 else DARK_GRAY
                paragraph.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Row coloring
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color
            elif i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = alt_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

    return table_shape


def fmt_num(n, decimals=0):
    """Format a number with commas."""
    if decimals == 0:
        return f"{int(round(n)):,}"
    return f"{n:,.{decimals}f}"


# ===================================================================
# SLIDE 1: Title
# ===================================================================
print("Creating Slide 1: Title")
s1 = add_slide()

# Full-slide dark blue background
bg = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BLUE
bg.line.fill.background()

# Green accent stripe
stripe = s1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), SLIDE_W, Inches(0.08)
)
stripe.fill.solid()
stripe.fill.fore_color.rgb = GREEN
stripe.line.fill.background()

# Title
add_textbox(s1, "CoastalCred", Inches(0.8), Inches(0.8),
            Inches(11.7), Inches(1.0), font_size=Pt(48), bold=True,
            color=WHITE, alignment=PP_ALIGN.LEFT)

add_textbox(s1, "Blockchain-Based Blue Carbon Registry & MRV System",
            Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.8),
            font_size=Pt(28), bold=False, color=WHITE, alignment=PP_ALIGN.LEFT)

# Subtitle info
add_textbox(s1, "Sem VI Mini Project  |  RCOEM, Dept. of Data Science, Nagpur",
            Inches(0.8), Inches(3.5), Inches(11.7), Inches(0.6),
            font_size=Pt(20), color=WHITE, alignment=PP_ALIGN.LEFT)

add_textbox(s1, "Team: Aniruddha Lahoti  |  Ansh Chopada  |  Tanmay Gaikwad  |  Sanika Bonde",
            Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.5),
            font_size=Pt(18), color=WHITE, alignment=PP_ALIGN.LEFT)

add_textbox(s1, "Guide: Dr. Aarti Karandikar   |   Industry Mentor: Rishikesh Kale (Filecoin/Protocol Labs)",
            Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.5),
            font_size=Pt(16), color=WHITE, alignment=PP_ALIGN.LEFT)


# ===================================================================
# SLIDE 2: Problem & Objectives
# ===================================================================
print("Creating Slide 2: Problem & Objectives")
s2 = add_slide()
add_header_bar(s2, "Problem & Objectives")
add_green_accent_line(s2)

# Left column - Problem
add_textbox(s2, "The Problem", Inches(0.7), Inches(1.3),
            Inches(5.5), Inches(0.5), font_size=Pt(22), bold=True, color=DARK_BLUE)

problem_bullets = [
    "Carbon markets lack transparency: voluntary credits face double-counting and fraud",
    "Manual MRV is slow & expensive: field surveys cost $50-100/ha, take months",
    "Mangroves sequester 3-5x more carbon than terrestrial forests, yet only 1% of blue carbon is credited",
]
add_bullet_box(s2, problem_bullets, Inches(0.7), Inches(1.9),
               Inches(5.5), Inches(4.5), font_size=Pt(16))

# Right column - Objectives
add_textbox(s2, "Our Objectives", Inches(7.0), Inches(1.3),
            Inches(5.5), Inches(0.5), font_size=Pt(22), bold=True, color=GREEN)

obj_bullets = [
    "ML pipeline for mangrove classification: satellite imagery to binary masks",
    "Compare 3 models: NDVI baseline, XGBoost, U-Net (rules to deep learning)",
    "IPCC Tier 1 carbon estimation: convert classified hectares to tCO\u2082e credits",
]
add_bullet_box(s2, obj_bullets, Inches(7.0), Inches(1.9),
               Inches(5.5), Inches(4.5), font_size=Pt(16))

# Divider line
div = s2.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(6.4), Inches(1.3), Inches(0.04), Inches(5.0)
)
div.fill.solid()
div.fill.fore_color.rgb = GREEN
div.line.fill.background()


# ===================================================================
# SLIDE 3: Data Pipeline & Study Sites
# ===================================================================
print("Creating Slide 3: Data Pipeline & Study Sites")
s3 = add_slide()
add_header_bar(s3, "Data Pipeline & Study Sites")
add_green_accent_line(s3)

data_bullets = [
    "3 Study Sites: Sundarbans, WB (Train) | Gulf of Kutch, GJ (Train) | Pichavaram, TN (Test)",
    "Satellite Data: Sentinel-2 L2A via Google Earth Engine, 6 bands (B2, B3, B4, B8, B11, B12)",
    "Temporal Scope: 2020 (baseline) and 2024 (current) -- two-point flux measurement",
    "Patch Extraction: 256x256 pixels, stride 128 (50% overlap)",
    "Site-Level Split: Train 6,374 | Val 708 | Test 71 patches",
    "Ground Truth: Global Mangrove Watch v3 polygons rasterized to Sentinel-2 grid",
]
add_bullet_box(s3, data_bullets, Inches(0.7), Inches(1.3),
               Inches(6.0), Inches(5.5), font_size=Pt(16))

# Embed alignment check image on right half
add_image(s3, "alignment_check_sundarbans_2024.png",
          Inches(7.0), Inches(1.3), width=Inches(5.8))


# ===================================================================
# SLIDE 4: Three Models Overview
# ===================================================================
print("Creating Slide 4: Three Models Overview")
s4 = add_slide()
add_header_bar(s4, "Three Models -- Complexity Ladder")
add_green_accent_line(s4)

# Sub-header
add_textbox(s4, "Rules  \u2192  Classical ML  \u2192  Deep Learning",
            Inches(0.7), Inches(1.25), Inches(12.0), Inches(0.5),
            font_size=Pt(20), bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)

# Three columns for three models
# Model 1: NDVI
ndvi_box = s4.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.0), Inches(3.8), Inches(4.5)
)
ndvi_box.fill.solid()
ndvi_box.fill.fore_color.rgb = LIGHT_BLUE_ROW
ndvi_box.line.color.rgb = DARK_BLUE
ndvi_box.line.width = Pt(2)

add_textbox(s4, "1. NDVI Threshold", Inches(0.7), Inches(2.2),
            Inches(3.4), Inches(0.5), font_size=Pt(22), bold=True, color=DARK_BLUE)

ndvi_bullets = [
    "Rule-based baseline",
    "NDVI = (B8 - B4) / (B8 + B4)",
    "Best threshold: 0.55 (grid search)",
    "No training required",
    "Establishes performance floor",
]
add_bullet_box(s4, ndvi_bullets, Inches(0.7), Inches(2.8),
               Inches(3.4), Inches(3.5), font_size=Pt(14), bold_prefix=False)

# Model 2: XGBoost
xgb_box = s4.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.7), Inches(2.0), Inches(3.8), Inches(4.5)
)
xgb_box.fill.solid()
xgb_box.fill.fore_color.rgb = LIGHT_BLUE_ROW
xgb_box.line.color.rgb = DARK_BLUE
xgb_box.line.width = Pt(2)

add_textbox(s4, "2. XGBoost", Inches(4.9), Inches(2.2),
            Inches(3.4), Inches(0.5), font_size=Pt(22), bold=True, color=DARK_BLUE)

xgb_bullets = [
    "Classical ML pixel classifier",
    "10 features: 6 bands + NDVI, EVI, NDWI, SAVI",
    "287 trees, scale_pos_weight=1.47",
    "Training time: 2.8 seconds",
    "Feature importances for interpretability",
]
add_bullet_box(s4, xgb_bullets, Inches(4.9), Inches(2.8),
               Inches(3.4), Inches(3.5), font_size=Pt(14), bold_prefix=False)

# Model 3: U-Net
unet_box = s4.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.9), Inches(2.0), Inches(3.8), Inches(4.5)
)
unet_box.fill.solid()
unet_box.fill.fore_color.rgb = LIGHT_BLUE_ROW
unet_box.line.color.rgb = DARK_BLUE
unet_box.line.width = Pt(2)

add_textbox(s4, "3. U-Net", Inches(9.1), Inches(2.2),
            Inches(3.4), Inches(0.5), font_size=Pt(22), bold=True, color=DARK_BLUE)

unet_bullets = [
    "Deep learning segmentation",
    "ResNet-18 encoder (ImageNet pretrained)",
    "6-channel input, 1 output class",
    "BCEWithLogitsLoss + pos_weight",
    "50 epochs, training: ~2 hours",
]
add_bullet_box(s4, unet_bullets, Inches(9.1), Inches(2.8),
               Inches(3.4), Inches(3.5), font_size=Pt(14), bold_prefix=False)


# ===================================================================
# SLIDE 5: Model Comparison -- Test Set Results (KEY SLIDE)
# ===================================================================
print("Creating Slide 5: Model Comparison (KEY SLIDE)")
s5 = add_slide()
add_header_bar(s5, "Model Comparison -- Test Set Results (Pichavaram)")
add_green_accent_line(s5)

# Comparison table
ndvi_test = ndvi_data["test_metrics"]
xgb_test = xgb_data["test_metrics"]
unet_test = unet_data["test_metrics"]

comp_rows = [
    ["Model", "Precision", "Recall", "IoU", "F1", "Training Time"],
    ["NDVI Threshold", f"{ndvi_test['precision']:.3f}",
     f"{ndvi_test['recall']:.3f}", f"{ndvi_test['iou']:.3f}",
     f"{ndvi_test['f1']:.3f}", "0s"],
    ["XGBoost", f"{xgb_test['precision']:.3f}",
     f"{xgb_test['recall']:.3f}", f"{xgb_test['iou']:.3f}",
     f"{xgb_test['f1']:.3f}", "2.8s"],
    ["U-Net", f"{unet_test['precision']:.3f}",
     f"{unet_test['recall']:.3f}", f"{unet_test['iou']:.3f}",
     f"{unet_test['f1']:.3f}", "7,299s"],
]

comp_col_widths = [Inches(2.2), Inches(1.5), Inches(1.5),
                   Inches(1.5), Inches(1.5), Inches(1.8)]
add_table(s5, comp_rows, comp_col_widths, Inches(0.7), Inches(1.3),
          height_per_row=Inches(0.55), font_size=Pt(16), header_font_size=Pt(17))

# Key insight callout
callout = s5.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(3.7),
    Inches(5.5), Inches(0.7)
)
callout.fill.solid()
callout.fill.fore_color.rgb = LIGHT_GREEN
callout.line.color.rgb = GREEN
callout.line.width = Pt(2)

add_textbox(s5, "Key Insight: NDVI baseline wins on test IoU (0.336) despite being the simplest model",
            Inches(0.9), Inches(3.8), Inches(5.1), Inches(0.5),
            font_size=Pt(15), bold=True, color=GREEN)

# Embed comparison chart
add_image(s5, "comparison_table.png",
          Inches(6.8), Inches(3.0), width=Inches(6.0))


# ===================================================================
# SLIDE 6: Generalization Gap & Per-Site Analysis (KEY SLIDE)
# ===================================================================
print("Creating Slide 6: Generalization Gap (KEY SLIDE)")
s6 = add_slide()
add_header_bar(s6, "Generalization Gap & Per-Site Analysis")
add_green_accent_line(s6)

# Generalization gap table
ndvi_val_iou = ndvi_data["val_metrics"]["iou"]
xgb_val_iou = xgb_data["val_metrics"]["iou"]
unet_val_iou = unet_data["val_metrics"]["iou"]

ndvi_test_iou = ndvi_data["test_metrics"]["iou"]
xgb_test_iou = xgb_data["test_metrics"]["iou"]
unet_test_iou = unet_data["test_metrics"]["iou"]

ndvi_drop = ((ndvi_test_iou - ndvi_val_iou) / ndvi_val_iou) * 100
xgb_drop = ((xgb_test_iou - xgb_val_iou) / xgb_val_iou) * 100
unet_drop = ((unet_test_iou - unet_val_iou) / unet_val_iou) * 100

gap_rows = [
    ["Model", "Val IoU", "Test IoU", "IoU Drop (%)"],
    ["NDVI Threshold", f"{ndvi_val_iou:.3f}", f"{ndvi_test_iou:.3f}", f"{ndvi_drop:.0f}%"],
    ["XGBoost", f"{xgb_val_iou:.3f}", f"{xgb_test_iou:.3f}", f"{xgb_drop:.0f}%"],
    ["U-Net", f"{unet_val_iou:.3f}", f"{unet_test_iou:.3f}", f"{unet_drop:.0f}%"],
]

gap_col_widths = [Inches(2.5), Inches(1.8), Inches(1.8), Inches(2.0)]
add_table(s6, gap_rows, gap_col_widths, Inches(0.7), Inches(1.3),
          height_per_row=Inches(0.55), font_size=Pt(16), header_font_size=Pt(17))

# Key insights on right side
insight_bullets = [
    "More complex models overfit more: U-Net drops 65%, NDVI drops only 41%",
    "Pichavaram is ecologically distinct: different species composition, tidal patterns",
    "Site-level split prevents data leakage: no shared patches between train and test",
    "Honest evaluation: real-world generalization, not inflated random-split metrics",
    "Takeaway: model complexity does not guarantee generalization",
]
add_bullet_box(s6, insight_bullets, Inches(0.7), Inches(3.8),
               Inches(12.0), Inches(3.0), font_size=Pt(16), bold_prefix=False)

# Embed Pichavaram alignment check on far right
add_image(s6, "alignment_check_pichavaram_2024.png",
          Inches(9.0), Inches(1.3), width=Inches(3.8))


# ===================================================================
# SLIDE 7: XGBoost Feature Importance
# ===================================================================
print("Creating Slide 7: XGBoost Feature Importance")
s7 = add_slide()
add_header_bar(s7, "XGBoost Feature Importance")
add_green_accent_line(s7)

# Large feature importance image
add_image(s7, "feature_importance.png",
          Inches(0.5), Inches(1.2), width=Inches(8.0))

# Caption and notes on right
fi_bullets = [
    "NDVI: 45.6% -- dominant feature (vegetation index)",
    "B2 Blue: 21.4% -- water/land discrimination",
    "B11 SWIR: 9.9% -- moisture content detection",
    "NDWI: 6.6% -- water index complements NDVI",
    "B8 NIR: 4.7% -- direct vegetation reflectance",
]
add_bullet_box(s7, fi_bullets, Inches(8.8), Inches(1.4),
               Inches(4.0), Inches(3.5), font_size=Pt(14))

# Insight box
add_textbox(s7, "Spectral indices (NDVI, NDWI, SAVI, EVI) contribute\n"
            "~54% of total importance -- engineered features\n"
            "outperform raw spectral bands.",
            Inches(8.8), Inches(5.0), Inches(4.0), Inches(1.8),
            font_size=Pt(14), bold=False, color=DARK_GRAY)


# ===================================================================
# SLIDE 8: U-Net Predictions
# ===================================================================
print("Creating Slide 8: U-Net Predictions")
s8 = add_slide()
add_header_bar(s8, "U-Net Prediction Visualizations")
add_green_accent_line(s8)

# Two prediction images stacked
add_image(s8, "unet_pred_0.png",
          Inches(0.3), Inches(1.2), width=Inches(6.2))

add_image(s8, "unet_pred_1.png",
          Inches(6.7), Inches(1.2), width=Inches(6.2))

# Caption
add_textbox(s8, "Sentinel-2 RGB  |  Ground Truth Mask  |  U-Net Prediction",
            Inches(0.3), Inches(6.8), Inches(12.7), Inches(0.5),
            font_size=Pt(16), bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

# Note box
note_box = s8.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(5.5),
    Inches(6.3), Inches(0.7)
)
note_box.fill.solid()
note_box.fill.fore_color.rgb = LIGHT_GREEN
note_box.line.color.rgb = GREEN
note_box.line.width = Pt(1)

add_textbox(s8, "High precision (0.908) but low recall (0.276) -- conservative predictions, few false positives",
            Inches(3.7), Inches(5.55), Inches(5.9), Inches(0.6),
            font_size=Pt(14), bold=True, color=GREEN)


# ===================================================================
# SLIDE 9: Carbon Credit Estimation
# ===================================================================
print("Creating Slide 9: Carbon Credit Estimation")
s9 = add_slide()
add_header_bar(s9, "Carbon Credit Estimation (IPCC Tier 1)")
add_green_accent_line(s9)

# Formula
add_textbox(s9, "IPCC Tier 1 Formula:",
            Inches(0.7), Inches(1.2), Inches(5.0), Inches(0.4),
            font_size=Pt(18), bold=True, color=DARK_BLUE)

add_textbox(s9, "Stock: hectares x 230 t/ha x 0.47 x 3.667 = tCO\u2082e\n"
            "Flux: (current_ha - baseline_ha) x 7.0 x years",
            Inches(0.7), Inches(1.65), Inches(6.0), Inches(0.8),
            font_size=Pt(15), bold=False, color=DARK_GRAY)

# Carbon flux table -- all 3 models x 3 sites
sites = ["sundarbans", "gulf_of_kutch", "pichavaram"]
site_labels = ["Sundarbans", "Gulf of Kutch", "Pichavaram"]

carbon_rows = [["Site", "NDVI Flux (4yr tCO\u2082e)", "XGBoost Flux (4yr tCO\u2082e)", "U-Net Flux (4yr tCO\u2082e)"]]

for site, label in zip(sites, site_labels):
    ndvi_flux = carbon_data["ndvi"][site]["carbon_flux"]["total_flux_tco2e_4yr"]
    xgb_flux = carbon_data["xgboost"][site]["carbon_flux"]["total_flux_tco2e_4yr"]
    unet_flux = carbon_data["unet"][site]["carbon_flux"]["total_flux_tco2e_4yr"]
    carbon_rows.append([
        label,
        fmt_num(ndvi_flux),
        fmt_num(xgb_flux),
        fmt_num(unet_flux),
    ])

carbon_col_widths = [Inches(2.2), Inches(2.8), Inches(2.8), Inches(2.8)]
add_table(s9, carbon_rows, carbon_col_widths, Inches(0.7), Inches(2.6),
          height_per_row=Inches(0.5), font_size=Pt(15), header_font_size=Pt(16))

# Key finding callout
cf_callout = s9.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(5.0),
    Inches(11.5), Inches(1.2)
)
cf_callout.fill.solid()
cf_callout.fill.fore_color.rgb = LIGHT_BLUE_ROW
cf_callout.line.color.rgb = DARK_BLUE
cf_callout.line.width = Pt(2)

carbon_insights = [
    "Models disagree on direction of change: NDVI shows Sundarbans gaining; XGBoost and U-Net show loss",
    "Carbon flux estimates vary by orders of magnitude depending on model accuracy",
    "Demonstrates why multi-model comparison is essential for credible MRV",
]
add_bullet_box(s9, carbon_insights, Inches(0.9), Inches(5.1),
               Inches(11.1), Inches(1.0), font_size=Pt(14), bold_prefix=False)


# ===================================================================
# SLIDE 10: Conclusion & Future Work
# ===================================================================
print("Creating Slide 10: Conclusion & Future Work")
s10 = add_slide()
add_header_bar(s10, "Conclusion & Future Work")
add_green_accent_line(s10)

# Left column - Key Findings
add_textbox(s10, "Key Findings", Inches(0.7), Inches(1.3),
            Inches(5.5), Inches(0.5), font_size=Pt(22), bold=True, color=DARK_BLUE)

findings = [
    "NDVI baseline competitive on unseen sites (IoU 0.336 vs XGBoost 0.332)",
    "Generalization is the main challenge: 41-65% IoU drop on test site",
    "U-Net achieves highest val performance (IoU 0.768) but overfits most",
    "Carbon flux estimates vary significantly by model -- consensus needed",
    "Site-level split provides honest, non-inflated evaluation",
]
add_bullet_box(s10, findings, Inches(0.7), Inches(1.9),
               Inches(5.5), Inches(4.0), font_size=Pt(15), bold_prefix=False)

# Right column - Sem VII Roadmap
add_textbox(s10, "Sem VII Roadmap", Inches(7.0), Inches(1.3),
            Inches(5.5), Inches(0.5), font_size=Pt(22), bold=True, color=GREEN)

roadmap = [
    "Smart contract deployment on Polygon Amoy testnet",
    "Mobile app for community field verification",
    "IPFS/Filecoin integration for immutable evidence storage",
    "Improved U-Net with data augmentation & domain adaptation",
    "Full microservices deployment with React frontend portals",
    "Integration with Verra/Gold Standard methodology",
]
add_bullet_box(s10, roadmap, Inches(7.0), Inches(1.9),
               Inches(5.5), Inches(4.0), font_size=Pt(15), bold_prefix=False)

# Divider
div2 = s10.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(6.4), Inches(1.3), Inches(0.04), Inches(5.0)
)
div2.fill.solid()
div2.fill.fore_color.rgb = GREEN
div2.line.fill.background()

# Thank You bar at bottom
ty_bar = s10.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.5), SLIDE_W, Inches(0.7)
)
ty_bar.fill.solid()
ty_bar.fill.fore_color.rgb = DARK_BLUE
ty_bar.line.fill.background()

add_textbox(s10, "Thank You -- Questions?",
            Inches(0.6), Inches(6.55), SLIDE_W - Inches(1.2), Inches(0.6),
            font_size=Pt(24), bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


# ===================================================================
# Save
# ===================================================================
prs.save(str(OUTPUT))
file_size = OUTPUT.stat().st_size
print(f"\nPresentation saved to: {OUTPUT}")
print(f"File size: {file_size:,} bytes ({file_size / 1024:.1f} KB)")
print(f"Total slides: {slide_number_counter[0]}")
